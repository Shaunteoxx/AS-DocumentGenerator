# CRD Generator API
from fastapi import FastAPI, HTTPException, UploadFile, File, Form, Cookie, Depends, Request, Response
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from pathlib import Path
import google.generativeai as genai
from docx import Document
from pptx import Presentation
import fitz
import io
import os
import re
import json
import logging
import datetime
import time
import asyncio
import httpx
import jwt
from jwt import PyJWKClient
from dotenv import load_dotenv

load_dotenv()

genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

app = FastAPI(title="CRD Generator API")

ALLOWED_ORIGINS = [
    "https://project-xwokx.vercel.app",
    "http://localhost:5173",
    "http://localhost:3000",
]

app.add_middleware(
    CORSMiddleware,
    allow_origins=ALLOWED_ORIGINS,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

CONTEXT_DIR = Path(__file__).parent / "crd_context_data"
CREDENTIALS_PATH = Path(__file__).parent / os.getenv("GOOGLE_CREDENTIALS_PATH", "credentials.json")
SHEET_ID = "1JOsrTwUMpJ9cKKXgAAdSZRd_mJquhajAAvnIsk6B__I"
SHEET_WORKSHEET = "Feature"
MODEL_NAME = "gemini-2.5-flash"

CORRIDOR_BASE_URL = os.getenv("CORRIDOR_BASE_URL", "https://www.corridor.cloud")
CORRIDOR_CLIENT_ID = os.getenv("CORRIDOR_CLIENT_ID", "")
CORRIDOR_API_KEY = os.getenv("CORRIDOR_CLOUD_API_KEY", "")
CORRIDOR_REDIRECT_URI = os.getenv("CORRIDOR_REDIRECT_URI", "")
DEV_AUTH_BYPASS = os.getenv("DEV_AUTH_BYPASS", "false").lower() == "true"

logger = logging.getLogger(__name__)

_jwks_client: PyJWKClient | None = None


def get_jwks_client() -> PyJWKClient:
    global _jwks_client
    if _jwks_client is None:
        _jwks_client = PyJWKClient(f"{CORRIDOR_BASE_URL}/api/oauth/jwks")
    return _jwks_client


def verify_token(token: str) -> dict:
    client = get_jwks_client()
    signing_key = client.get_signing_key_from_jwt(token)
    return jwt.decode(token, signing_key.key, algorithms=["RS256"], options={"verify_aud": False})


async def require_auth(
    request: Request,
    access_token: str | None = Cookie(default=None),
) -> dict:
    if DEV_AUTH_BYPASS:
        return {"sub": "dev"}
    # Prefer Authorization header (Bearer token from sessionStorage)
    auth_header = request.headers.get("Authorization", "")
    token = auth_header.removeprefix("Bearer ").strip() if auth_header.startswith("Bearer ") else access_token
    if not token:
        logger.warning("require_auth: no token present")
        raise HTTPException(status_code=401, detail="Not authenticated")
    try:
        # verify_token can hit the network (JWKS refresh) and does CPU-bound
        # signature checks — keep it off the event loop.
        return await asyncio.to_thread(verify_token, token)
    except jwt.ExpiredSignatureError:
        logger.warning("require_auth: token expired")
        raise HTTPException(status_code=401, detail="Token expired")
    except Exception as e:
        logger.warning("require_auth: token validation failed: %s", e)
        raise HTTPException(status_code=401, detail="Invalid token")


def parse_questions(raw: str) -> list[str]:
    """Best-effort extraction of a JSON array of question strings from model output.

    Handles code fences (```json …```), surrounding prose, and responses wrapped
    as an object like {"questions": [...]}.
    """
    text = (raw or "").strip()
    if "```" in text:
        m = re.search(r"```(?:json)?\s*(.*?)```", text, re.DOTALL)
        if m:
            text = m.group(1).strip()

    candidates = [text]
    start, end = text.find("["), text.rfind("]")
    if start != -1 and end > start:
        candidates.append(text[start:end + 1])

    for candidate in candidates:
        try:
            data = json.loads(candidate)
        except (json.JSONDecodeError, TypeError):
            continue
        if isinstance(data, dict):
            data = data.get("questions") or next((v for v in data.values() if isinstance(v, list)), None)
        if isinstance(data, list):
            questions = [str(q).strip() for q in data if str(q).strip()]
            if questions:
                return questions

    logger.error("parse_questions failed | raw: %s", (raw or "")[:300])
    raise HTTPException(status_code=500, detail="Failed to parse questions from model response")


def parse_json_response(raw: str, error_detail: str) -> dict:
    """Parse a JSON object from model output, tolerating ``` fences and prose."""
    text = (raw or "").strip()
    if "```" in text:
        text = text.split("```")[1]
        if text.startswith("json"):
            text = text[4:]
        text = text.strip()
    text = text.rstrip("`").strip()
    try:
        return json.loads(text)
    except json.JSONDecodeError as exc:
        logger.error("%s: %s | raw: %s", error_detail, exc, (raw or "")[:300])
        raise HTTPException(status_code=500, detail=error_detail)


def strip_code_fences(text: str) -> str:
    """Strip a wrapping ``` / ```markdown fence the model sometimes puts around
    the entire document. Left in, that wrapper makes the whole doc render as one
    literal code block — raw #, **, - shown as text — when the markdown is
    converted to a Google Doc on export."""
    s = (text or "").strip()
    if not s.startswith("```"):
        return s
    lines = s.split("\n")
    lines = lines[1:]  # drop the opening fence line (``` + optional language tag)
    if lines and lines[-1].strip().startswith("```"):
        lines = lines[:-1]  # drop the closing fence line
    return "\n".join(lines).strip()


_features_cache: dict = {"data": None, "ts": 0.0}
FEATURES_CACHE_TTL = 300  # seconds


def get_features_from_sheet() -> str:
    now = time.time()
    if _features_cache["data"] is not None and now - _features_cache["ts"] < FEATURES_CACHE_TTL:
        return _features_cache["data"]
    try:
        import gspread
        from google.oauth2.service_account import Credentials

        scopes = [
            "https://www.googleapis.com/auth/spreadsheets.readonly",
            "https://www.googleapis.com/auth/drive.readonly",
        ]
        creds = Credentials.from_service_account_file(str(CREDENTIALS_PATH), scopes=scopes)
        gc = gspread.authorize(creds)
        ws = gc.open_by_key(SHEET_ID).worksheet(SHEET_WORKSHEET)
        rows = ws.get_all_records(expected_headers=['ID', 'Domain', 'Feature Name', 'Key User Story', 'Status', 'User Persona', 'Dependencies', 'BrdID'])

        current_statuses = {"Completed", "In Progress"}
        current, future = [], []
        for row in rows:
            status = str(row.get("Status", "")).strip()
            entry = (
                f"  ID: {row.get('ID', '')} | Domain: {row.get('Domain', '')} | "
                f"Feature: {row.get('Feature Name', '')} | "
                f"User Story: {row.get('Key User Story', '')} | "
                f"Status: {status} | Persona: {row.get('User Persona', '')} | "
                f"Dependencies: {row.get('Dependencies', '')} | BRD: {row.get('BrdID', '')}"
            )
            if status in current_statuses:
                current.append(entry)
            else:
                future.append(entry)

        parts = []
        if current:
            parts.append("=== CURRENT FEATURES ===\n" + "\n".join(current))
        if future:
            parts.append("=== FUTURE FEATURES ===\n" + "\n".join(future))
        result = "\n\n".join(parts) if parts else "No features found in sheet."
        _features_cache["data"] = result
        _features_cache["ts"] = now
        return result

    except Exception as exc:
        logger.warning("Google Sheets fetch failed (%s), falling back to local files.", exc)
        fallback_parts = []
        for name in ("allo8_current.txt", "allo8_future.txt"):
            p = CONTEXT_DIR / name
            if p.exists():
                fallback_parts.append(p.read_text())
        return "\n\n".join(fallback_parts) if fallback_parts else ""


def extract_client_name_from_crd(md: str) -> str:
    match = re.search(r'\|\s*\*{0,2}Client\s+Name\*{0,2}\s*\|\s*([^|\n]+)\|', md, re.IGNORECASE)
    if match:
        return match.group(1).strip()
    return ""


def _load_context_dir(directory: Path, empty_msg: str) -> str:
    parts = []
    for f in sorted(directory.glob("*")):
        if f.is_file() and f.suffix in {".txt", ".md"}:
            parts.append(f"### {f.name}\n{f.read_text()}")
    return "\n\n".join(parts) if parts else empty_msg


def _make_model(role: str, context: str) -> genai.GenerativeModel:
    system = f"""{role}

Always follow the corporate context and templates below when generating documents.

--- CORPORATE CONTEXT ---
{context}
--- END CORPORATE CONTEXT ---"""
    return genai.GenerativeModel(model_name=MODEL_NAME, system_instruction=system)


def load_context() -> str:
    return _load_context_dir(CONTEXT_DIR, "No corporate context loaded.")


def get_model(context: str) -> genai.GenerativeModel:
    return _make_model(
        "You are a CRD (Customer Request Document) generator assistant. You help analysts create professional CRDs by analyzing client notes, asking targeted clarifying questions, and producing formatted CRD documents.",
        context,
    )


SUPPORTED_EXTENSIONS = {".txt", ".md", ".docx", ".pdf", ".pptx"}


def extract_text(ext: str, content: bytes) -> str:
    if ext in {".txt", ".md"}:
        return content.decode("utf-8", errors="replace")
    elif ext == ".docx":
        doc = Document(io.BytesIO(content))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    elif ext == ".pdf":
        with fitz.open(stream=content, filetype="pdf") as pdf:
            return "\n".join(page.get_text() for page in pdf)
    elif ext == ".pptx":
        prs = Presentation(io.BytesIO(content))
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text)
        return "\n".join(parts)
    return ""


def _name_to_crd_id(client_name: str, version: int = 1) -> str:
    words = [w for w in client_name.split() if w]
    if len(words) >= 2:
        initials = words[0][0].upper() + words[1][0].upper()
    elif len(words) == 1 and len(words[0]) >= 2:
        initials = words[0][:2].upper()
    elif len(words) == 1:
        initials = (words[0][0] * 2).upper()
    else:
        return f"C-XX-{version}"
    return f"C-{initials}-{version}"


def generate_crd_id(md: str, version: int = 1) -> str:
    match = re.search(r'\|\s*\*{0,2}Client\s+Name\*{0,2}\s*\|\s*([^|\n]+)\|', md, re.IGNORECASE)
    if not match:
        return f"C-XX-{version}"
    return _name_to_crd_id(match.group(1).strip(), version)


def extract_id_from_notes(text: str, version: int = 1) -> str:
    patterns = [
        r'(?:Client|Company|Account|Customer)(?:\s+Name)?\s*[:\-]\s*([^\n,\.]{2,80})',
        r'\|\s*\*{0,2}Client\s+Name\*{0,2}\s*\|\s*([^|\n]+)\|',
    ]
    for pat in patterns:
        m = re.search(pat, text, re.IGNORECASE)
        if m:
            name = m.group(1).strip().rstrip('.,')
            if name:
                return _name_to_crd_id(name, version)
    return f"C-XX-{version}"


# ── Auth endpoints ────────────────────────────────────────────────────────────

class TokenExchangeRequest(BaseModel):
    code: str
    code_verifier: str


@app.post("/auth/token")
async def auth_token(req: TokenExchangeRequest, response: Response):
    async with httpx.AsyncClient() as client:
        r = await client.post(
            f"{CORRIDOR_BASE_URL}/api/oauth/token",
            json={
                "grant_type": "authorization_code",
                "client_id": CORRIDOR_CLIENT_ID,
                "code": req.code,
                "redirect_uri": CORRIDOR_REDIRECT_URI,
                "code_verifier": req.code_verifier,
            },
            headers={"X-API-Key": CORRIDOR_API_KEY},
        )
    if r.status_code != 200:
        logger.error("Corridor token exchange failed: %s %s", r.status_code, r.text)
        raise HTTPException(status_code=400, detail="Token exchange failed")
    data = r.json()
    is_secure = not DEV_AUTH_BYPASS
    cookie_opts = dict(httponly=True, secure=is_secure, samesite="none" if is_secure else "lax")
    response.set_cookie("access_token", data["access_token"], max_age=data.get("expires_in", 3600), **cookie_opts)
    response.set_cookie("refresh_token", data["refresh_token"], max_age=30 * 24 * 3600, **cookie_opts)
    return {"ok": True, "access_token": data["access_token"], "expires_in": data.get("expires_in", 3600)}


@app.post("/auth/refresh")
async def auth_refresh(response: Response, refresh_token: str | None = Cookie(default=None)):
    if not refresh_token:
        raise HTTPException(status_code=401, detail="No refresh token")
    async with httpx.AsyncClient() as client:
        r = await client.post(
            f"{CORRIDOR_BASE_URL}/api/oauth/token",
            json={
                "grant_type": "refresh_token",
                "client_id": CORRIDOR_CLIENT_ID,
                "refresh_token": refresh_token,
            },
            headers={"X-API-Key": CORRIDOR_API_KEY},
        )
    if r.status_code != 200:
        raise HTTPException(status_code=401, detail="Token refresh failed")
    data = r.json()
    is_secure = not DEV_AUTH_BYPASS
    cookie_opts = dict(httponly=True, secure=is_secure, samesite="none" if is_secure else "lax")
    response.set_cookie("access_token", data["access_token"], max_age=data.get("expires_in", 3600), **cookie_opts)
    if "refresh_token" in data:
        response.set_cookie("refresh_token", data["refresh_token"], max_age=30 * 24 * 3600, **cookie_opts)
    return {"ok": True, "access_token": data["access_token"], "expires_in": data.get("expires_in", 3600)}


@app.post("/auth/logout")
async def auth_logout(response: Response):
    response.delete_cookie("access_token", samesite="none", secure=True)
    response.delete_cookie("refresh_token", samesite="none", secure=True)
    return {"ok": True}


@app.get("/auth/me")
async def auth_me(claims: dict = Depends(require_auth)):
    return claims


# ── Existing endpoints (all protected) ───────────────────────────────────────

@app.post("/upload")
async def upload_file(
    file: UploadFile = File(...),
    _: dict = Depends(require_auth),
):
    ext = Path(file.filename).suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        raise HTTPException(status_code=400, detail=f"Unsupported file type: {ext}")
    content = await file.read()
    return {"text": extract_text(ext, content)}


class ClarifyRequest(BaseModel):
    notes: str
    analysis: str
    answers: list[dict] = []


class GenerateRequest(BaseModel):
    notes: str
    analysis: str
    answers: list[dict]
    filename: str = ""




class RegenerateRequest(BaseModel):
    crd: str
    section: str = ""
    instruction: str = ""


class LogToSheetRequest(BaseModel):
    filename: str
    client_name: str
    drive_link: str = ""


@app.post("/analyze")
async def analyze(
    notes: str = Form(default=""),
    files: list[UploadFile] = File(default=[]),
    _: dict = Depends(require_auth),
):
    file_parts = []
    for file in files:
        ext = Path(file.filename).suffix.lower()
        if ext not in SUPPORTED_EXTENSIONS:
            continue
        text = extract_text(ext, await file.read())
        if text.strip():
            file_parts.append(f"--- {file.filename} ---\n{text}")

    combined = "\n\n".join(filter(None, [notes.strip()] + file_parts))
    if not combined.strip():
        raise HTTPException(status_code=400, detail="No content provided")

    model = get_model(load_context())
    response = await model.generate_content_async(
        f"""Analyze these client notes and extract:
1. Key requirements
2. Stakeholders mentioned
3. Scope and deliverables
4. Any ambiguities or missing information

Client Notes:
{combined}"""
    )
    return {"analysis": response.text, "combined_notes": combined}


@app.post("/clarify")
async def clarify(req: ClarifyRequest, _: dict = Depends(require_auth)):
    model = get_model(load_context())
    answers_text = ""
    if req.answers:
        answers_text = "\n\nPrevious Q&A:\n" + "\n".join(
            f"Q: {a['question']}\nA: {a['answer']}" for a in req.answers
        )

    features_text = await asyncio.to_thread(get_features_from_sheet)
    features_block = (
        f"\n\n--- ALLOCATE SPACE PLATFORM FEATURES (live from product sheet) ---\n{features_text}\n---"
        if features_text else ""
    )

    response = await model.generate_content_async(
        f"""Based on the client notes and analysis below, generate 3–5 targeted clarifying questions needed to complete a CRD. Only ask what is truly necessary. Do not ask about platform features or capabilities — these are already known from the product sheet above.

Client Notes:
{req.notes}

Analysis:
{req.analysis}{answers_text}{features_block}

Respond with ONLY a JSON array of question strings, no other text. Example: ["Question 1?", "Question 2?"]"""
    )

    questions = parse_questions(response.text)

    return {"questions": questions}


@app.post("/generate")
async def generate(req: GenerateRequest, _: dict = Depends(require_auth)):
    model = get_model(load_context())
    answers_text = "\n".join(
        f"Q: {a['question']}\nA: {a['answer']}" for a in req.answers
    )

    filename = req.filename.strip()

    if not filename:
        filename = extract_id_from_notes(req.notes + "\n" + req.analysis)

    filename_instruction = (
        f"\n\nThe Docs ID for this document is: {filename}. "
        "Use this exact value in the Docs ID field of the Document Control section. "
        "Do not generate a new ID."
    )

    features_text = await asyncio.to_thread(get_features_from_sheet)
    features_block = (
        f"\n\n--- ALLOCATE SPACE PLATFORM FEATURES (live from product sheet) ---\n{features_text}\n---"
        if features_text else ""
    )

    today = datetime.date.today().isoformat()
    response = await model.generate_content_async(
        f"""Generate a complete, professionally formatted CRD document following the corporate template in your context.

Client Notes:
{req.notes}

Analysis:
{req.analysis}

Clarifying Q&A:
{answers_text}{features_block}

Today's date is {today}. Use this for the Date Prepared field.
Produce the full CRD in Markdown format.{filename_instruction}"""
    )

    crd = strip_code_fences(response.text)
    crd_id = filename

    crd = re.sub(
        r'(\|\s*\*{0,2}Docs\s+ID\*{0,2}\s*\|)([^|\n]*)(\|)',
        lambda m: f"{m.group(1)} {crd_id} {m.group(3)}",
        crd, count=1, flags=re.IGNORECASE,
    )

    return {"crd": crd, "crd_id": crd_id}


@app.post("/regenerate")
async def regenerate_section(req: RegenerateRequest, _: dict = Depends(require_auth)):
    if not req.section.strip():
        raise HTTPException(status_code=400, detail="section is required")
    model = get_model(load_context())
    instruction_line = f"\nAdditional instructions: {req.instruction.strip()}" if req.instruction.strip() else ""
    response = await model.generate_content_async(
        f"""You are given a complete CRD document in Markdown format. Rewrite ONLY the section named below. Do not alter any other section, heading, or content.

Section to regenerate: {req.section.strip()}{instruction_line}

Full CRD:
{req.crd}

Return ONLY the complete updated Markdown document with that section rewritten. No preamble, no explanation."""
    )
    return {"crd": strip_code_fences(response.text)}



def _open_worksheet(worksheet: str | int):
    """Sync gspread helper — always call via asyncio.to_thread."""
    import gspread
    from google.oauth2.service_account import Credentials

    scopes = ["https://www.googleapis.com/auth/spreadsheets"]
    creds = Credentials.from_service_account_file(str(CREDENTIALS_PATH), scopes=scopes)
    gc = gspread.authorize(creds)
    sheet = gc.open_by_key(SHEET_ID)
    if isinstance(worksheet, int):
        return sheet.get_worksheet_by_id(worksheet)
    return sheet.worksheet(worksheet)


@app.post("/log-to-sheet")
async def log_to_sheet(req: LogToSheetRequest, _: dict = Depends(require_auth)):
    date_str = datetime.date.today().strftime("%-d %b %y")

    def _write():
        ws = _open_worksheet("CR")
        ws.append_row(
            [req.filename, req.client_name, "", req.drive_link, date_str, "", "", ""],
            value_input_option="RAW",
            insert_data_option="INSERT_ROWS",
        )

    try:
        await asyncio.to_thread(_write)
    except Exception as exc:
        logger.warning("CR sheet write failed in /log-to-sheet (%s)", exc)
        return {"status": "warning", "message": str(exc)}
    return {"status": "ok"}


# ── IRD helpers ───────────────────────────────────────────────────────────────

IRD_CONTEXT_DIR = Path(__file__).parent / "ird_context_data"


def load_ird_context() -> str:
    return _load_context_dir(IRD_CONTEXT_DIR, "No IRD context loaded.")


def get_ird_model(context: str) -> genai.GenerativeModel:
    return _make_model(
        "You are an IRD (Internal Requirement Document) generator assistant for Allocate Space. You help teams document internal operational requirements by analyzing source documents, asking targeted clarifying questions, and producing formatted IRD documents.",
        context,
    )


def _name_to_ird_id(name: str, version: int = 1) -> str:
    words = [w for w in name.split() if w]
    if len(words) >= 2:
        initials = words[0][0].upper() + words[1][0].upper()
    elif len(words) == 1 and len(words[0]) >= 2:
        initials = words[0][:2].upper()
    elif len(words) == 1:
        initials = (words[0][0] * 2).upper()
    else:
        return f"I-XX-{version}"
    return f"I-{initials}-{version}"


def extract_ird_id(text: str, version: int = 1) -> str:
    patterns = [
        r'(?:Team|Department|Group|Division)(?:\s+Name)?\s*[:\-]\s*([^\n,\.]{2,80})',
        r'(?:Requestor|Raised By|Submitted By)\s*[:\-]\s*([^\n,\.]{2,80})',
    ]
    for pat in patterns:
        m = re.search(pat, text, re.IGNORECASE)
        if m:
            name = m.group(1).strip().rstrip('.,')
            if name:
                return _name_to_ird_id(name, version)
    return f"I-XX-{version}"


class IrdRegenerateRequest(BaseModel):
    crd: str
    section: str = ""
    instruction: str = ""


@app.post("/ird/analyze")
async def ird_analyze(
    notes: str = Form(default=""),
    files: list[UploadFile] = File(default=[]),
    _: dict = Depends(require_auth),
):
    file_parts = []
    for file in files:
        ext = Path(file.filename).suffix.lower()
        if ext not in SUPPORTED_EXTENSIONS:
            continue
        text = extract_text(ext, await file.read())
        if text.strip():
            file_parts.append(f"--- {file.filename} ---\n{text}")

    combined = "\n\n".join(filter(None, [notes.strip()] + file_parts))
    if not combined.strip():
        raise HTTPException(status_code=400, detail="No content provided")

    model = get_ird_model(load_ird_context())
    response = await model.generate_content_async(
        f"""Analyze these internal documents and extract:
1. Key internal requirements
2. Teams or stakeholders mentioned
3. Scope and operational needs
4. Any ambiguities or missing information

Source Documents:
{combined}"""
    )
    return {"analysis": response.text, "combined_notes": combined}


@app.post("/ird/clarify")
async def ird_clarify(req: ClarifyRequest, _: dict = Depends(require_auth)):
    model = get_ird_model(load_ird_context())
    answers_text = ""
    if req.answers:
        answers_text = "\n\nPrevious Q&A:\n" + "\n".join(
            f"Q: {a['question']}\nA: {a['answer']}" for a in req.answers
        )

    response = await model.generate_content_async(
        f"""Based on the internal documents and analysis below, generate 3–5 targeted clarifying questions needed to complete an IRD. Only ask what is truly necessary.

Source Documents:
{req.notes}

Analysis:
{req.analysis}{answers_text}

Respond with ONLY a JSON array of question strings, no other text. Example: ["Question 1?", "Question 2?"]"""
    )

    questions = parse_questions(response.text)

    return {"questions": questions}


@app.post("/ird/generate")
async def ird_generate(req: GenerateRequest, _: dict = Depends(require_auth)):
    model = get_ird_model(load_ird_context())
    answers_text = "\n".join(
        f"Q: {a['question']}\nA: {a['answer']}" for a in req.answers
    )

    filename = req.filename.strip()
    if not filename:
        filename = extract_ird_id(req.notes + "\n" + req.analysis)

    filename_instruction = (
        f"\n\nThe Docs ID for this document is: {filename}. "
        "Use this exact value in the Docs ID field of the Document Control section. "
        "Do not generate a new ID."
    )

    today = datetime.date.today().isoformat()
    response = await model.generate_content_async(
        f"""Generate a complete, professionally formatted IRD document following the corporate template in your context.

Source Documents:
{req.notes}

Analysis:
{req.analysis}

Clarifying Q&A:
{answers_text}

Today's date is {today}. Use this for the Date Prepared field.
Produce the full IRD in Markdown format.{filename_instruction}"""
    )

    ird = strip_code_fences(response.text)
    return {"crd": ird, "crd_id": filename}


@app.post("/ird/regenerate")
async def ird_regenerate(req: IrdRegenerateRequest, _: dict = Depends(require_auth)):
    if not req.section.strip():
        raise HTTPException(status_code=400, detail="section is required")
    model = get_ird_model(load_ird_context())
    instruction_line = f"\nAdditional instructions: {req.instruction.strip()}" if req.instruction.strip() else ""
    response = await model.generate_content_async(
        f"""You are given a complete IRD document in Markdown format. Rewrite ONLY the section named below. Do not alter any other section, heading, or content.

Section to regenerate: {req.section.strip()}{instruction_line}

Full IRD:
{req.crd}

Return ONLY the complete updated Markdown document with that section rewritten. No preamble, no explanation."""
    )
    return {"crd": strip_code_fences(response.text)}


class IrdLogToSheetRequest(BaseModel):
    ird: str
    filename: str
    drive_link: str = ""


def extract_ird_initiative_name(md: str) -> str:
    match = re.search(
        r'\*\*(?:Initiative\s+Name|Team\s*/\s*Department)\*\*[:\s]+([^\n*]+)',
        md, re.IGNORECASE,
    )
    if match:
        return match[1].strip()
    heading = re.search(r'^#\s+(.+)$', md, re.MULTILINE)
    return heading.group(1).strip() if heading else "Internal Requirement"


def extract_ird_key_request(md: str) -> str:
    match = re.search(r'-\s*\*\*Key Request:\*\*\s*([^\n]+)', md, re.IGNORECASE)
    if match:
        return match.group(1).strip().strip('_()')
    return ""


@app.post("/ird/log-to-sheet")
async def ird_log_to_sheet(req: IrdLogToSheetRequest, _: dict = Depends(require_auth)):
    initiative_name = extract_ird_initiative_name(req.ird)
    description = extract_ird_key_request(req.ird)
    doc_link = f'=HYPERLINK("{req.drive_link}","{initiative_name}")' if req.drive_link else initiative_name

    def _write():
        ws = _open_worksheet(1456248746)
        ws.append_row(
            [req.filename, initiative_name, description, "", doc_link],
            value_input_option="USER_ENTERED",
            insert_data_option="INSERT_ROWS",
        )

    try:
        await asyncio.to_thread(_write)
    except Exception as exc:
        logger.warning("IR sheet write failed in /ird/log-to-sheet (%s)", exc)
        return {"status": "warning", "message": str(exc)}
    return {"status": "ok"}


# ── PRD helpers ───────────────────────────────────────────────────────────────

PRD_CONTEXT_DIR = Path(__file__).parent / "prd_context_data"


def load_prd_context() -> str:
    return _load_context_dir(PRD_CONTEXT_DIR, "No PRD context loaded.")


def get_prd_model(context: str) -> genai.GenerativeModel:
    return _make_model(
        "You are a PRD (Product Requirement Document) generator assistant for Allocate Space. You help product managers create professional PRDs by analyzing source documents, asking targeted clarifying questions, and producing formatted PRD documents.",
        context,
    )


class PrdRegenerateRequest(BaseModel):
    crd: str
    section: str = ""
    instruction: str = ""


@app.post("/prd/analyze")
async def prd_analyze(
    notes: str = Form(default=""),
    files: list[UploadFile] = File(default=[]),
    _: dict = Depends(require_auth),
):
    file_parts = []
    for file in files:
        ext = Path(file.filename).suffix.lower()
        if ext not in SUPPORTED_EXTENSIONS:
            continue
        text = extract_text(ext, await file.read())
        if text.strip():
            file_parts.append(f"--- {file.filename} ---\n{text}")

    combined = "\n\n".join(filter(None, [notes.strip()] + file_parts))
    if not combined.strip():
        raise HTTPException(status_code=400, detail="No content provided")

    model = get_prd_model(load_prd_context())
    response = await model.generate_content_async(
        f"""Analyze these product documents and extract:
1. Key product requirements and user needs
2. Target user personas mentioned
3. Scope and feature boundaries
4. Any ambiguities or missing information

Source Documents:
{combined}"""
    )
    return {"analysis": response.text, "combined_notes": combined}


@app.post("/prd/clarify")
async def prd_clarify(req: ClarifyRequest, _: dict = Depends(require_auth)):
    model = get_prd_model(load_prd_context())
    answers_text = ""
    if req.answers:
        answers_text = "\n\nPrevious Q&A:\n" + "\n".join(
            f"Q: {a['question']}\nA: {a['answer']}" for a in req.answers
        )

    response = await model.generate_content_async(
        f"""Based on the source documents and analysis below, generate 3–5 targeted clarifying questions needed to complete a PRD. Only ask what is truly necessary.

Source Documents:
{req.notes}

Analysis:
{req.analysis}{answers_text}

Respond with ONLY a JSON array of question strings, no other text. Example: ["Question 1?", "Question 2?"]"""
    )

    questions = parse_questions(response.text)

    return {"questions": questions}


def extract_prd_title(prd: str) -> str:
    """Extract the H1 title from a generated PRD (e.g. 'FBK1 - Form Block V1.0 PRD')."""
    m = re.search(r'^#\s+(.+)$', prd, re.MULTILINE)
    return m.group(1).strip() if m else 'prd'


@app.post("/prd/generate")
async def prd_generate(req: GenerateRequest, _: dict = Depends(require_auth)):
    model = get_prd_model(load_prd_context())
    answers_text = "\n".join(
        f"Q: {a['question']}\nA: {a['answer']}" for a in req.answers
    )

    today = datetime.date.today().isoformat()
    response = await model.generate_content_async(
        f"""Generate a complete, professionally formatted PRD document following the corporate template in your context.

Source Documents:
{req.notes}

Analysis:
{req.analysis}

Clarifying Q&A:
{answers_text}

Today's date is {today}. Use this for the Date Prepared field.
Produce the full PRD in Markdown format."""
    )

    prd = strip_code_fences(response.text)
    filename = extract_prd_title(prd)
    return {"crd": prd, "crd_id": filename}


@app.post("/prd/regenerate")
async def prd_regenerate(req: PrdRegenerateRequest, _: dict = Depends(require_auth)):
    if not req.section.strip():
        raise HTTPException(status_code=400, detail="section is required")
    model = get_prd_model(load_prd_context())
    instruction_line = f"\nAdditional instructions: {req.instruction.strip()}" if req.instruction.strip() else ""
    response = await model.generate_content_async(
        f"""You are given a complete PRD document in Markdown format. Rewrite ONLY the section named below. Do not alter any other section, heading, or content.

Section to regenerate: {req.section.strip()}{instruction_line}

Full PRD:
{req.crd}

Return ONLY the complete updated Markdown document with that section rewritten. No preamble, no explanation."""
    )
    return {"crd": strip_code_fences(response.text)}


CLICKUP_API_KEY = os.getenv("CLICKUP_API_KEY", "")
CLICKUP_WORKSPACE_ID = "9008246823"
CLICKUP_PRD_FOLDER_ID = "901814228652"


def extract_prd_feature_name(md: str) -> str:
    match = re.search(r'\*\*Feature\s+Name\*\*[:\s]+([^\n*]+)', md, re.IGNORECASE)
    if match:
        return match[1].strip()
    heading = re.search(r'^#\s+(.+)$', md, re.MULTILINE)
    return heading.group(1).strip() if heading else "Product Requirement Document"


class PrdExportRequest(BaseModel):
    prd: str
    filename: str = ""


@app.post("/prd/export-to-clickup")
async def prd_export_to_clickup(req: PrdExportRequest, _: dict = Depends(require_auth)):
    if not CLICKUP_API_KEY:
        raise HTTPException(status_code=500, detail="ClickUp API key not configured")

    title = req.filename.strip() or extract_prd_feature_name(req.prd)
    headers = {"Authorization": CLICKUP_API_KEY, "Content-Type": "application/json"}

    async with httpx.AsyncClient(timeout=30) as client:
        create_res = await client.post(
            f"https://api.clickup.com/api/v3/workspaces/{CLICKUP_WORKSPACE_ID}/docs",
            headers=headers,
            json={"name": title, "parent": {"id": CLICKUP_PRD_FOLDER_ID, "type": 5}},
        )
        if create_res.status_code not in (200, 201):
            raise HTTPException(status_code=502, detail=f"ClickUp doc creation failed: {create_res.text}")

        doc_data = create_res.json()
        doc_id = doc_data.get("id") or (doc_data.get("doc") or {}).get("id")
        if not doc_id:
            raise HTTPException(status_code=502, detail="ClickUp did not return a doc ID")

        page_res = await client.post(
            f"https://api.clickup.com/api/v3/workspaces/{CLICKUP_WORKSPACE_ID}/docs/{doc_id}/pages",
            headers=headers,
            json={"name": title, "content": req.prd, "content_format": "text/md"},
        )
        if page_res.status_code not in (200, 201):
            logger.warning("ClickUp page creation failed (%s): %s", page_res.status_code, page_res.text)

    doc_url = doc_data.get("url") or f"https://app.clickup.com/{CLICKUP_WORKSPACE_ID}/v/dc/{doc_id}"
    return {"doc_url": doc_url, "doc_id": doc_id}


# ── BRD helpers ───────────────────────────────────────────────────────────────

BRD_CONTEXT_DIR = Path(__file__).parent / "brd_context_data"


def load_brd_context() -> str:
    return _load_context_dir(BRD_CONTEXT_DIR, "No BRD context loaded.")


def get_brd_model(context: str) -> genai.GenerativeModel:
    return _make_model(
        "You are a BRD (Business Requirements Document) generator assistant for Allocate Space. You help analysts create professional BRDs by analyzing source CRD documents, asking targeted clarifying questions, and producing formatted BRD documents.",
        context,
    )


def extract_brd_title(md: str) -> str:
    match = re.search(r'^#\s+(.+)$', md, re.MULTILINE)
    if match:
        return match.group(1).strip()
    return "Business Requirement"


def generate_brd_filename(md: str) -> str:
    title = extract_brd_title(md)
    today = datetime.date.today()
    date_str = today.strftime("%b") + str(today.day)
    return f"{title} - {date_str} - BRD"


def generate_br_initials(title: str) -> str:
    clean = re.sub(r'[^a-zA-Z0-9\s]', ' ', title)
    words = [w for w in clean.split() if w.isalpha()]
    return ''.join(w[0].upper() for w in words) or "BR"


def extract_brd_objective(md: str) -> str:
    match = re.search(
        r'\*\*Key Solution Objective\*\*\s*[:\-]?\s*[\*_]?([^\n\*_]+)[\*_]?',
        md, re.IGNORECASE,
    )
    if match:
        return match.group(1).strip()
    return ""


def extract_brd_crds(md: str) -> str:
    # CR IDs use the form C-GTS / C-SPEDG (no trailing number), but also accept
    # an optional requirement suffix like C-GTS-04. Collect from the "Originating
    # CR(s)" header field and the Existing Reference Material section.
    regions = []
    orig = re.search(r'\*\*Originating CR\(s\):?\*\*\s*([^\n]+)', md, re.IGNORECASE)
    if orig:
        regions.append(orig.group(1))
    section = re.search(
        r'##\s*Existing Reference Material\s*\n(.*?)(?=\n##|\Z)',
        md, re.DOTALL | re.IGNORECASE,
    )
    if section:
        regions.append(section.group(1))

    crd_ids = []
    for region in regions:
        crd_ids.extend(re.findall(r'\bC-[A-Z]+(?:-\d+)?\b', region))
    return ", ".join(dict.fromkeys(crd_ids))


# ── BRD request models ────────────────────────────────────────────────────────

class BrdRegenerateRequest(BaseModel):
    brd: str
    section: str = ""
    instruction: str = ""


class BrdLogToSheetRequest(BaseModel):
    brd: str
    drive_link: str = ""


# ── BRD Drive folder ──────────────────────────────────────────────────────────

BRD_DRIVE_FOLDER_ID = "1F2_IRbCAwltGPI1i4UaSgpGDtjVxWsfx"


_brd_corpus_cache: dict = {"data": None, "ts": 0.0}
BRD_CORPUS_CACHE_TTL = 300  # seconds


async def _fetch_one_brd(client: httpx.AsyncClient, token: str, f: dict) -> dict | None:
    """Download and extract the text of a single Drive BRD file."""
    file_id = f["id"]
    mime = f.get("mimeType", "")
    name = f.get("name", "")
    link = f.get("webViewLink", f"https://drive.google.com/file/d/{file_id}/view")

    # Skip template/config files so the AI only matches against real BRDs.
    name_lower = name.lower()
    if (
        re.search(r'(^|\s|-)template(\s|-|$)', name_lower) is not None
        or name_lower.endswith((".md", ".txt"))
        or "application/vnd.google-apps.shortcut" in mime
        or "application/vnd.google-apps.folder" in mime
    ):
        logger.info("BRD fetch: skipping template/shortcut file '%s'", name)
        return None

    try:
        if "google-apps.document" in mime:
            er = await client.get(
                f"https://www.googleapis.com/drive/v3/files/{file_id}/export",
                headers={"Authorization": f"Bearer {token}"},
                params={"mimeType": "text/plain"},
                timeout=15,
            )
            text = er.text if er.is_success else ""
        else:
            dr = await client.get(
                f"https://www.googleapis.com/drive/v3/files/{file_id}?alt=media",
                headers={"Authorization": f"Bearer {token}"},
                timeout=15,
            )
            text = extract_text(Path(name).suffix.lower(), dr.content) if dr.is_success else ""

        if text.strip():
            return {"name": name, "id": file_id, "link": link, "text": text[:10000]}
    except Exception as exc:
        logger.warning("Failed to fetch BRD '%s': %s", name, exc)
    return None


async def fetch_brd_texts_from_drive(refresh: bool = False) -> list[dict]:
    """Return [{name, id, link, text}] for every BRD in the Drive folder.

    Results are cached for BRD_CORPUS_CACHE_TTL seconds since the BRD corpus
    rarely changes between generations; downloads run concurrently.
    """
    now = time.time()
    if not refresh and _brd_corpus_cache["data"] is not None \
            and now - _brd_corpus_cache["ts"] < BRD_CORPUS_CACHE_TTL:
        return _brd_corpus_cache["data"]

    try:
        token = await _get_drive_token(readonly=True)
    except Exception as exc:
        logger.warning("Drive token fetch failed: %s", exc)
        return []

    results = []
    try:
        async with httpx.AsyncClient(timeout=30) as client:
            r = await client.get(
                "https://www.googleapis.com/drive/v3/files",
                headers={"Authorization": f"Bearer {token}"},
                params={
                    "q": f"'{BRD_DRIVE_FOLDER_ID}' in parents and trashed=false",
                    "fields": "files(id,name,webViewLink,mimeType)",
                    "pageSize": 100,
                },
            )
            if not r.is_success:
                logger.warning("Drive file list failed: %s", r.text)
                return []

            fetched = await asyncio.gather(
                *(_fetch_one_brd(client, token, f) for f in r.json().get("files", []))
            )
            results = [b for b in fetched if b is not None]
    except Exception as exc:
        logger.warning("Drive folder fetch error: %s", exc)
        return _brd_corpus_cache["data"] or []

    _brd_corpus_cache["data"] = results
    _brd_corpus_cache["ts"] = now
    return results


# ── Document graph ────────────────────────────────────────────────────────────

CRD_DRIVE_FOLDER_ID = "1MTojq7o5eU6ypCb7JrXhnpo4Q34X8UzF"
IRD_DRIVE_FOLDER_ID = "10SMxLiD4paaAtP2QFu-XAywsct9L8t-0"
PRD_DRIVE_FOLDER_ID = "1UrlSNK_-6BOBVl9RHqRbWIEFbl_VMtju"

_graph_cache: dict = {"data": None, "ts": 0.0}
GRAPH_CACHE_TTL = 300  # seconds


async def _fetch_one_folder_file(client: httpx.AsyncClient, token: str, f: dict, doc_type: str) -> dict | None:
    """Download and extract the text of a single Drive file for the graph."""
    file_id, mime, name = f["id"], f.get("mimeType", ""), f.get("name", "")
    link = f.get("webViewLink", f"https://drive.google.com/file/d/{file_id}/view")
    try:
        if "google-apps.document" in mime:
            er = await client.get(
                f"https://www.googleapis.com/drive/v3/files/{file_id}/export",
                headers={"Authorization": f"Bearer {token}"},
                params={"mimeType": "text/plain"}, timeout=15,
            )
            text = er.text if er.is_success else ""
        else:
            dr = await client.get(
                f"https://www.googleapis.com/drive/v3/files/{file_id}?alt=media",
                headers={"Authorization": f"Bearer {token}"}, timeout=15,
            )
            text = extract_text(Path(name).suffix.lower(), dr.content) if dr.is_success else ""
        return {"id": file_id, "name": name, "type": doc_type,
                "link": link, "modified": f.get("modifiedTime", ""), "text": text[:8000]}
    except Exception as exc:
        logger.warning("Failed to read %s '%s': %s", doc_type, name, exc)
        return None


async def _fetch_folder(folder_id: str, doc_type: str) -> list[dict]:
    """Generic version of fetch_brd_texts_from_drive for any folder."""
    try:
        token = await _get_drive_token(readonly=True)
    except Exception as exc:
        logger.warning("Drive token failed for %s: %s", doc_type, exc)
        return []

    try:
        async with httpx.AsyncClient(timeout=30) as client:
            r = await client.get(
                "https://www.googleapis.com/drive/v3/files",
                headers={"Authorization": f"Bearer {token}"},
                params={
                    "q": f"'{folder_id}' in parents and trashed=false",
                    "fields": "files(id,name,webViewLink,mimeType,modifiedTime)",
                    "pageSize": 100,
                },
            )
            docs = []
            for f in r.json().get("files", []):
                name_lower, mime = f.get("name", "").lower(), f.get("mimeType", "")
                # Only filter explicit template/config files, not any doc that happens to have "template" in the name
                is_template = (
                    re.search(r'(^|\s|-)template(\s|-|$)', name_lower) is not None
                    or name_lower.endswith((".md", ".txt"))
                    or "application/vnd.google-apps.shortcut" in mime
                    or "application/vnd.google-apps.folder" in mime
                )
                if is_template:
                    logger.info("Graph: skipping template/shortcut file '%s' (%s)", f.get("name", ""), doc_type)
                    continue
                docs.append(f)

            fetched = await asyncio.gather(
                *(_fetch_one_folder_file(client, token, f, doc_type) for f in docs)
            )
            return [d for d in fetched if d is not None]
    except Exception as exc:
        logger.warning("Folder fetch failed for %s: %s", doc_type, exc)
        return []


def _re_field(text: str, field: str) -> str:
    # Match both **Field**: value (markdown) and Field: value (plain-text export)
    m = re.search(rf'(?:\*\*{re.escape(field)}\*\*|{re.escape(field)}):?\s*([^\n\r]+)', text, re.IGNORECASE)
    if not m:
        return ""
    v = re.sub(r'<!--.*?-->', '', m.group(1)).strip()
    v = re.sub(r'^[\*_(]+|[\*_)]+$', '', v).strip()
    if not v:
        return ""
    bad_starts = ('_', 'Leave blank', 'To be', '] —', '[Version', '[Feature', '[Product')
    return "" if any(v.startswith(p) for p in bad_starts) else v


def _drive_id_from_url(url: str) -> str:
    m = re.search(r'/d/([a-zA-Z0-9_-]{20,})', url)
    return m.group(1) if m else ""


def _parse_node(raw: dict) -> dict:
    text, t = raw["text"], raw["type"]
    if t == "crd":
        docs_id = _re_field(text, "Docs ID")
        if not docs_id:
            # Old-style CRD header: "C-XXX-v2.0-1  Client Request Document"
            hm = re.search(r'(C-[A-Z0-9]+-[v.\d-]*\d)\s+Client Request Document', text, re.IGNORECASE)
            docs_id = hm.group(1).strip() if hm else raw["name"]
        title = _re_field(text, "Client Name") or raw["name"]
        linked_brd = _re_field(text, "Linked BRD")
        source_ref = ""
    elif t == "ird":
        docs_id = _re_field(text, "Docs ID") or raw["name"]
        title = _re_field(text, "Initiative Name") or _re_field(text, "Domain") or raw["name"]
        linked_brd = _re_field(text, "Linked BRD")
        source_ref = ""
    elif t == "brd":
        m = re.search(r'^#\s+(.+)$', text, re.MULTILINE)
        docs_id = raw["name"]
        title = m.group(1).strip() if m else raw["name"]
        linked_brd = ""
        source_ref = _re_field(text, "Initial CR")
    elif t == "prd":
        docs_id = _re_field(text, "Docs ID") or raw["name"]
        title = _re_field(text, "Feature Name") or raw["name"]
        linked_brd = ""
        source_ref = ""
    else:
        docs_id = title = raw["name"]
        linked_brd = source_ref = ""
    return {
        "id": raw["id"], "docs_id": docs_id, "type": t, "title": title,
        "drive_link": raw["link"], "modified": raw["modified"][:10] if raw["modified"] else "",
        "_linked_brd": linked_brd, "_source_ref": source_ref, "_text": text,
    }


def _build_explicit_edges(nodes: list[dict]) -> list[dict]:
    by_id = {n["id"]: n for n in nodes}
    brds = [n for n in nodes if n["type"] == "brd"]
    crds_irds = [n for n in nodes if n["type"] in ("crd", "ird")]
    edges, seen = [], set()

    def add(src, tgt, method):
        k = (src, tgt)
        if k not in seen and src in by_id and tgt in by_id:
            seen.add(k)
            edges.append({"source": src, "target": tgt, "method": method})

    for n in crds_irds:
        ref = n["_linked_brd"]
        if not ref:
            continue
        drive_id = _drive_id_from_url(ref)
        if drive_id and drive_id in by_id:
            add(n["id"], drive_id, "explicit")
        else:
            for brd in brds:
                if ref.lower() in brd["title"].lower() or ref.lower() in brd["docs_id"].lower():
                    add(n["id"], brd["id"], "explicit")
                    break

    for brd in brds:
        ref = brd["_source_ref"]
        if not ref:
            continue
        ref_lower = ref.lower()
        for ci in crds_irds:
            # Use word-boundary match on docs_id to avoid false substring matches
            # (e.g. "Gim Tian Civil Engineering" title appearing inside another CRD's description)
            did = re.escape(ci["docs_id"].lower())
            if re.search(rf'(?<![a-z0-9]){did}(?![a-z0-9])', ref_lower):
                add(ci["id"], brd["id"], "explicit")

    return edges


async def _infer_edges(nodes: list[dict], explicit_edges: list[dict]) -> list[dict]:
    already_linked = {e["source"] for e in explicit_edges}
    unlinked_sources = [n for n in nodes if n["type"] in ("crd", "ird") and n["id"] not in already_linked]
    brds = [n for n in nodes if n["type"] == "brd"]
    prds = [n for n in nodes if n["type"] == "prd"]
    valid_ids = {n["id"] for n in nodes}
    inferred = []

    model = genai.GenerativeModel(MODEL_NAME)

    def _summarise(node_list):
        return "\n".join(
            f'- ID: {n["id"]} | {n["type"].upper()}: {n["title"]} | {n["_text"][:400]}'
            for n in node_list
        )

    async def _infer(prompt: str, label: str) -> list[dict]:
        edges = []
        try:
            resp = await model.generate_content_async(prompt)
            m = re.search(r'\{.*\}', resp.text, re.DOTALL)
            if m:
                for e in json.loads(m.group()).get("edges", []):
                    if e.get("confidence", 0) >= 0.5 and e["source"] in valid_ids and e["target"] in valid_ids:
                        edges.append({"source": e["source"], "target": e["target"],
                                      "method": "inferred", "confidence": round(e["confidence"], 2)})
        except Exception as exc:
            logger.warning("AI inference (%s) failed: %s", label, exc)
        return edges

    tasks = []
    if unlinked_sources and brds:
        tasks.append(_infer(
            f"""Determine which CRD/IRD documents are related to which BRD documents based on content alignment.

CRD/IRD (sources):
{_summarise(unlinked_sources)}

BRD (targets):
{_summarise(brds)}

Return ONLY JSON: {{"edges":[{{"source":"drive_id","target":"drive_id","confidence":0.9}}]}}
Only create edges with clear content alignment. confidence 0.5–1.0. Empty array if none.""",
            "CRD→BRD",
        ))
    if brds and prds:
        tasks.append(_infer(
            f"""Determine which BRD documents are fulfilled by which PRD documents.

BRD (sources — defines business needs):
{_summarise(brds)}

PRD (targets — defines product features being built):
{_summarise(prds)}

Return ONLY JSON: {{"edges":[{{"source":"brd_drive_id","target":"prd_drive_id","confidence":0.9}}]}}
Only create edges where the PRD clearly addresses the BRD's requirements. confidence 0.5–1.0.""",
            "BRD→PRD",
        ))

    for batch in await asyncio.gather(*tasks):
        inferred.extend(batch)

    return inferred


@app.get("/graph")
async def get_document_graph(refresh: bool = False, _: dict = Depends(require_auth)):
    global _graph_cache
    now = time.time()
    if not refresh and _graph_cache["data"] and now - _graph_cache["ts"] < GRAPH_CACHE_TTL:
        return {**_graph_cache["data"], "cached": True}

    raw_all = await asyncio.gather(
        _fetch_folder(CRD_DRIVE_FOLDER_ID, "crd"),
        _fetch_folder(IRD_DRIVE_FOLDER_ID, "ird"),
        _fetch_folder(BRD_DRIVE_FOLDER_ID, "brd"),
        _fetch_folder(PRD_DRIVE_FOLDER_ID, "prd"),
    )
    nodes = [_parse_node(r) for batch in raw_all for r in batch]

    explicit = _build_explicit_edges(nodes)
    inferred = await _infer_edges(nodes, explicit)

    clean_nodes = [{k: v for k, v in n.items() if not k.startswith("_")} for n in nodes]
    data = {
        "nodes": clean_nodes,
        "edges": explicit + inferred,
        "fetched_at": datetime.datetime.now().isoformat(),
        "cached": False,
    }
    _graph_cache = {"data": data, "ts": now}
    return data


@app.get("/graph/debug")
async def get_graph_debug(_: dict = Depends(require_auth)):
    """Returns raw parsed fields for each document — useful for diagnosing missing links."""
    raw_all = await asyncio.gather(
        _fetch_folder(CRD_DRIVE_FOLDER_ID, "crd"),
        _fetch_folder(IRD_DRIVE_FOLDER_ID, "ird"),
        _fetch_folder(BRD_DRIVE_FOLDER_ID, "brd"),
        _fetch_folder(PRD_DRIVE_FOLDER_ID, "prd"),
    )
    nodes = [_parse_node(r) for batch in raw_all for r in batch]
    explicit = _build_explicit_edges(nodes)
    debug = []
    for n in nodes:
        debug.append({
            "id": n["id"],
            "type": n["type"],
            "name_from_drive": next((r["name"] for batch in raw_all for r in batch if r["id"] == n["id"]), ""),
            "title": n["title"],
            "docs_id": n["docs_id"],
            "linked_brd_raw": n["_linked_brd"],
            "source_ref_raw": n["_source_ref"],
            "has_explicit_edge": any(e["source"] == n["id"] or e["target"] == n["id"] for e in explicit),
            "text_snippet": n["_text"][:300],
        })
    return {"nodes": debug, "explicit_edges": explicit}


# ── BRD endpoints ─────────────────────────────────────────────────────────────

@app.post("/brd/refresh-corpus")
async def brd_refresh_corpus(_: dict = Depends(require_auth)):
    """Force a re-fetch of the BRD corpus from Drive, bypassing the TTL cache."""
    brds = await fetch_brd_texts_from_drive(refresh=True)
    return {"count": len(brds), "names": [b["name"] for b in brds]}


@app.post("/brd/analyze")
async def brd_analyze(
    notes: str = Form(default=""),
    files: list[UploadFile] = File(default=[]),
    _: dict = Depends(require_auth),
):
    file_parts = []
    for file in files:
        ext = Path(file.filename).suffix.lower()
        if ext not in SUPPORTED_EXTENSIONS:
            continue
        text = extract_text(ext, await file.read())
        if text.strip():
            file_parts.append(f"--- {file.filename} ---\n{text}")

    combined = "\n\n".join(filter(None, [notes.strip()] + file_parts))
    if not combined.strip():
        raise HTTPException(status_code=400, detail="No content provided")

    existing_brds = await fetch_brd_texts_from_drive()
    brd_blocks = (
        "\n\n".join(
            f"=== {b['name']} ===\nDrive link: {b['link']}\n\n{b['text']}"
            for b in existing_brds
        )
        if existing_brds
        else "No existing BRDs found."
    )

    model = get_brd_model(load_brd_context())
    response = await model.generate_content_async(
        f"""Analyze the uploaded CRD documents and compare each distinct client requirement against the existing BRDs listed below.

EXTRACTION RULES — follow these exactly:
- Extract requirements only from the Business Requirements section of the CRD. This is the section containing numbered items labelled BR-01, BR-02, BR-03, etc.
- Do not extract requirements from the Underlying Problems, Desired Outcome, Trigger Event, Request Summary, or any other section.
- Each numbered BR item in the Business Requirements section is exactly one requirement. Extract precisely those items — no more, no fewer.
- The requirement name must come from the Business Name field of each BR item, not from the description, problem statement, or any other field.
- Capture the BR number (e.g. BR-01, BR-02) for each item as the br_id field. If no BR number can be identified, set br_id to null.
- Do not group, summarise, combine, or infer — extract faithfully from what is written in the Business Requirements section.
- If the CRD does not have a clearly labelled Business Requirements section with numbered BR items, only then fall back to inferring requirements from the document as a whole.

For each extracted requirement, determine one of three match states:
- matched: an existing BRD substantially and comprehensively covers the ENTIRE requirement. Be conservative — only use this if every distinct capability in the requirement is covered.
- partial: an existing BRD covers some aspects but there are clear gaps or new aspects not captured.
- unmatched: no existing BRD covers this requirement at all.

IMPORTANT: Many requirements describe several distinct capabilities (e.g. "manage users, configure solutions, and set up workspaces"). If an existing BRD covers only some of those capabilities, the status is partial — NOT matched — even if the part it covers is covered well. Only use matched when nothing of substance is left uncovered.

For partial matches, the coverage_note must explain what is already covered AND, specifically, which distinct capabilities are still missing.

CRITICAL — matched_brds:
- List EVERY existing BRD that covers any part of the requirement in the matched_brds array, each with its exact name and Drive link from the EXISTING BRDs list below. A requirement is often covered by several BRDs — include all of them, not just one.
- Use only documents that appear in the EXISTING BRDs list. Never invent a BRD name or link.
- For unmatched requirements, matched_brds must be an empty array [].
- Base your coverage judgement ONLY on the actual text of the EXISTING BRDs provided below. Do not claim a capability is missing if it is present in any of the provided BRDs.

Return ONLY a valid JSON object — no preamble, no explanation, no markdown fences:
{{
  "requirements": [
    {{
      "name": "Short requirement name",
      "description": "One sentence describing this requirement from the CRD.",
      "crd_source": "C-XX-01",
      "br_id": "BR-01",
      "status": "matched",
      "matched_brds": [
        {{ "name": "Exact BRD document name", "link": "https://drive.google.com/..." }}
      ],
      "coverage_note": null
    }},
    {{
      "name": "Another requirement",
      "description": "One sentence.",
      "crd_source": "C-XX-01",
      "br_id": "BR-02",
      "status": "partial",
      "matched_brds": [
        {{ "name": "First covering BRD", "link": "https://drive.google.com/..." }},
        {{ "name": "Second covering BRD", "link": "https://drive.google.com/..." }}
      ],
      "coverage_note": "These BRDs cover X and Y but do not capture Z."
    }},
    {{
      "name": "New requirement",
      "description": "One sentence.",
      "crd_source": "C-XX-01",
      "br_id": "BR-03",
      "status": "unmatched",
      "matched_brds": [],
      "coverage_note": null
    }}
  ]
}}

UPLOADED CRDs:
{combined}

EXISTING BRDs:
{brd_blocks}"""
    )

    data = parse_json_response(response.text, "Failed to parse match results from model")

    # Guard against hallucinated matches: every BRD a requirement claims to be
    # covered by must actually exist in the Drive folder. Validate the whole
    # matched_brds list; if nothing valid remains, downgrade to unmatched.
    def _file_id(link: str) -> str:
        m = re.search(r'/d/([a-zA-Z0-9_-]+)', link or "")
        if not m:
            m = re.search(r'[?&]id=([a-zA-Z0-9_-]+)', link or "")
        return m.group(1) if m else ""

    by_name = {b["name"].strip().lower(): b for b in existing_brds}
    by_id = {_file_id(b["link"]): b for b in existing_brds if _file_id(b["link"])}

    requirements = data.get("requirements", [])
    for req in requirements:
        # Normalise: accept the new array, or fall back to a single legacy field.
        claimed = req.get("matched_brds")
        if not isinstance(claimed, list):
            if req.get("matched_brd"):
                claimed = [{"name": req.get("matched_brd"), "link": req.get("matched_brd_link")}]
            else:
                claimed = []

        # Keep only entries that resolve to a real BRD, de-duplicated.
        valid, seen = [], set()
        for entry in claimed:
            if not isinstance(entry, dict):
                continue
            real = by_name.get((entry.get("name") or "").strip().lower()) or by_id.get(_file_id(entry.get("link") or ""))
            if real and real["id"] not in seen:
                seen.add(real["id"])
                valid.append({"name": real["name"], "link": real["link"]})

        if req.get("status") in ("matched", "partial") and not valid:
            logger.warning("Downgrading hallucinated match: req=%r claimed=%r", req.get("name"), claimed)
            req["status"] = "unmatched"
            req["coverage_note"] = None

        req["matched_brds"] = valid
        # Primary match kept for backward compatibility (e.g. the gap/update flow).
        req["matched_brd"] = valid[0]["name"] if valid else None
        req["matched_brd_link"] = valid[0]["link"] if valid else None

    return {"requirements": requirements, "combined_notes": combined}


class AnalyzeGapRequest(BaseModel):
    matched_brd_link: str
    requirement_name: str
    requirement_description: str
    crd_source: str
    coverage_note: str = ""


@app.post("/brd/analyze-gap")
async def brd_analyze_gap(req: AnalyzeGapRequest, _: dict = Depends(require_auth)):
    m = re.search(r'/d/([a-zA-Z0-9_-]+)', req.matched_brd_link)
    if not m:
        # Fall back: try query param id=
        m = re.search(r'[?&]id=([a-zA-Z0-9_-]+)', req.matched_brd_link)
    if not m:
        raise HTTPException(status_code=400, detail="Cannot parse file ID from Drive link")
    file_id = m.group(1)

    try:
        token = await _get_drive_token(readonly=True)
    except Exception as exc:
        raise HTTPException(status_code=500, detail=f"Drive auth failed: {exc}")

    async with httpx.AsyncClient(timeout=30) as client:
        er = await client.get(
            f"https://www.googleapis.com/drive/v3/files/{file_id}/export",
            headers={"Authorization": f"Bearer {token}"},
            params={"mimeType": "text/plain"},
        )
        if er.is_success:
            existing_brd_text = er.text
        else:
            dr = await client.get(
                f"https://www.googleapis.com/drive/v3/files/{file_id}?alt=media",
                headers={"Authorization": f"Bearer {token}"},
            )
            existing_brd_text = dr.text if dr.is_success else ""

    if not existing_brd_text.strip():
        raise HTTPException(status_code=500, detail="Could not retrieve existing BRD content from Drive")

    model = get_brd_model(load_brd_context())
    response = await model.generate_content_async(
        f"""Analyze the gap between an existing BRD and a new client requirement. Identify exactly what needs to change in the existing BRD to cover this new requirement.

New Requirement: {req.requirement_name}
Description: {req.requirement_description}
Source CRD: {req.crd_source}
Known gap: {req.coverage_note}

Existing BRD:
{existing_brd_text[:6000]}

Return ONLY a valid JSON object — no preamble, no explanation, no markdown fences:
{{
  "summary": "One sentence describing the overall gap",
  "sections": [
    {{
      "section": "Exact section name from the BRD",
      "change_type": "add",
      "reason": "Why this section needs to change",
      "draft": "The exact new content to add — ready to paste into the document"
    }}
  ]
}}"""
    )

    return parse_json_response(response.text, "Failed to parse gap report from model")


@app.post("/brd/clarify")
async def brd_clarify(req: ClarifyRequest, _: dict = Depends(require_auth)):
    model = get_brd_model(load_brd_context())
    answers_text = ""
    if req.answers:
        answers_text = "\n\nPrevious Q&A:\n" + "\n".join(
            f"Q: {a['question']}\nA: {a['answer']}" for a in req.answers
        )

    response = await model.generate_content_async(
        f"""Based on the source documents and analysis below, generate 3–5 targeted clarifying questions needed to define the scope and structure of the BRD.

Focus questions on:
- Scope boundaries (what is explicitly in or out of scope for this BRD)
- Dependencies between requirements across different CRDs or systems
- Ambiguities or contradictions in the source documents
- Priority or sequencing where requirements conflict or overlap

Source Documents:
{req.notes}

Analysis:
{req.analysis}{answers_text}

Respond with ONLY a JSON array of question strings, no other text. Example: ["Question 1?", "Question 2?"]"""
    )

    questions = parse_questions(response.text)

    return {"questions": questions}


@app.post("/brd/generate")
async def brd_generate(req: GenerateRequest, _: dict = Depends(require_auth)):
    model = get_brd_model(load_brd_context())
    answers_text = "\n".join(
        f"Q: {a['question']}\nA: {a['answer']}" for a in req.answers
    )

    today = datetime.date.today().isoformat()

    response = await model.generate_content_async(
        f"""Generate a complete, professionally formatted BRD document following the BRD template in your context.

Source Documents:
{req.notes}

Analysis:
{req.analysis}

Clarifying Q&A:
{answers_text}

Use today's date ({today}) for First Created.
Produce the full BRD in Markdown format."""
    )

    brd = strip_code_fences(response.text)
    brd_id = generate_brd_filename(brd)
    return {"brd": brd, "brd_id": brd_id}


@app.post("/brd/regenerate")
async def brd_regenerate(req: BrdRegenerateRequest, _: dict = Depends(require_auth)):
    if not req.section.strip():
        raise HTTPException(status_code=400, detail="section is required")
    model = get_brd_model(load_brd_context())
    instruction_line = f"\nAdditional instructions: {req.instruction.strip()}" if req.instruction.strip() else ""
    response = await model.generate_content_async(
        f"""You are given a complete BRD document in Markdown format. Rewrite ONLY the section named below. Do not alter any other section, heading, or content.

Section to regenerate: {req.section.strip()}{instruction_line}

Full BRD:
{req.brd}

Return ONLY the complete updated Markdown document with that section rewritten. No preamble, no explanation."""
    )
    return {"brd": strip_code_fences(response.text)}


# ── Document Review Hub ───────────────────────────────────────────────────────

DOCS_FOLDER_MAP = {
    "crd": CRD_DRIVE_FOLDER_ID,
    "brd": BRD_DRIVE_FOLDER_ID,
    "ird": IRD_DRIVE_FOLDER_ID,
}


async def _get_drive_token(readonly: bool = True) -> str:
    def _fetch() -> str:
        from google.oauth2.service_account import Credentials
        from google.auth.transport.requests import Request as GoogleAuthRequest
        scope = (
            "https://www.googleapis.com/auth/drive.readonly"
            if readonly
            else "https://www.googleapis.com/auth/drive"
        )
        creds = Credentials.from_service_account_file(str(CREDENTIALS_PATH), scopes=[scope])
        creds.refresh(GoogleAuthRequest())  # sync HTTP — keep off the event loop
        return creds.token

    return await asyncio.to_thread(_fetch)


async def _list_folder_metadata(client: httpx.AsyncClient, token: str, folder_id: str, doc_type: str) -> list[dict]:
    r = await client.get(
        "https://www.googleapis.com/drive/v3/files",
        headers={"Authorization": f"Bearer {token}"},
        params={
            "q": f"'{folder_id}' in parents and trashed=false",
            "fields": "files(id,name,webViewLink,mimeType,modifiedTime)",
            "pageSize": 100,
        },
    )
    if not r.is_success:
        return []
    docs = []
    for f in r.json().get("files", []):
        name_lower = f.get("name", "").lower()
        mime = f.get("mimeType", "")
        if (
            re.search(r'(^|\s|-)template(\s|-|$)', name_lower)
            or name_lower.endswith((".md", ".txt"))
            or "shortcut" in mime
            or "folder" in mime
        ):
            continue
        docs.append({
            "id": f["id"],
            "name": f["name"],
            "type": doc_type,
            "drive_link": f.get("webViewLink", f"https://drive.google.com/file/d/{f['id']}/view"),
            "modified": f.get("modifiedTime", "")[:10],
        })
    return docs


@app.get("/documents")
async def list_documents(_: dict = Depends(require_auth)):
    token = await _get_drive_token(readonly=True)
    async with httpx.AsyncClient(timeout=30) as client:
        results = await asyncio.gather(*(
            _list_folder_metadata(client, token, folder_id, doc_type)
            for doc_type, folder_id in DOCS_FOLDER_MAP.items()
        ))
    docs = [doc for batch in results for doc in batch]
    docs.sort(key=lambda d: d["modified"], reverse=True)
    return {"documents": docs}


def _docx_to_markdown(docx_bytes: bytes) -> str:
    """Convert DOCX bytes to markdown-structured text, preserving heading levels.

    DOCX export from Google Docs reliably maps Heading 1 / Heading 2 / Heading 3
    styles to proper paragraph styles regardless of how the doc was originally
    created — unlike HTML export, which can use class-based styles instead of
    semantic <h1>/<h2> tags when the doc was uploaded as HTML.
    """
    doc = Document(io.BytesIO(docx_bytes))

    def _heading_level(style_name: str) -> int:
        m = re.match(r'Heading (\d)', style_name or '')
        return int(m.group(1)) if m else 0

    def _is_list_item(para) -> bool:
        pPr = para._p.pPr
        return pPr is not None and pPr.numPr is not None

    def _para_to_md(para) -> str:
        level = _heading_level(para.style.name if para.style else '')
        # Headings: keep text clean (no ** wrapping) so it matches the Doc's heading
        # text exactly during surgical section replacement.
        if level:
            text = ''.join(r.text for r in para.runs).strip()
            return ('#' * level + ' ' + text) if text else ''
        # Body: preserve bold runs and bullet-list markers.
        text = ''.join(
            (f'**{r.text}**' if r.bold and r.text.strip() else r.text)
            for r in para.runs
        ).strip()
        if not text:
            return ''
        if _is_list_item(para):
            return '- ' + text
        return text

    def _table_to_md(table) -> list[str]:
        rows = []
        for i, row in enumerate(table.rows):
            cells = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
            rows.append('| ' + ' | '.join(cells) + ' |')
            if i == 0:
                rows.append('|' + '|'.join([' --- ' for _ in cells]) + '|')
        return rows

    from docx.text.paragraph import Paragraph as DocxParagraph
    from docx.table import Table as DocxTable

    parts = []
    # Iterate body children to preserve paragraph/table order
    for child in doc.element.body:
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag == 'p':
            md = _para_to_md(DocxParagraph(child, doc))
            if md:
                parts.append(md)
        elif tag == 'tbl':
            parts.extend(_table_to_md(DocxTable(child, doc)))

    text = '\n\n'.join(parts)
    text = re.sub(r'\n{3,}', '\n\n', text)
    return text.strip()


def _parse_sections(text: str) -> list[dict]:
    """Split a markdown-structured document into sections at every #, ## or ### heading.

    Splits on H1, H2 and H3 so docs work whether their section titles are styled
    as Heading 1, 2 or 3 (e.g. each BR-0N requirement is an H3).  The first chunk
    (document title / preamble before the first heading) is included only if it
    has non-trivial content beyond just the H1 title line.
    """
    parts = re.split(r'(?=\n#{1,3} )', text)
    sections = []
    for i, part in enumerate(parts):
        trimmed = part.strip()
        if not trimmed:
            continue
        first_line = trimmed.split('\n')[0]
        m = re.match(r'^#{1,3}\s+(.+)', first_line)
        heading = m.group(1).strip() if m else "Document Overview"

        # Skip the preamble if it is *only* the H1 title line (nothing else to show)
        rest = '\n'.join(trimmed.split('\n')[1:]).strip()
        if i == 0 and not rest and m and m.group(0).startswith('#') and not m.group(0).startswith('##'):
            continue

        sections.append({"heading": heading, "content": trimmed})
    return sections


@app.get("/documents/{doc_id}/content")
async def get_document_content(doc_id: str, doc_type: str = "crd", _: dict = Depends(require_auth)):
    token = await _get_drive_token(readonly=True)
    async with httpx.AsyncClient(timeout=30) as client:
        meta_r = await client.get(
            f"https://www.googleapis.com/drive/v3/files/{doc_id}",
            headers={"Authorization": f"Bearer {token}"},
            params={"fields": "id,name,mimeType,webViewLink"},
        )
        if not meta_r.is_success:
            raise HTTPException(status_code=404, detail="Document not found")
        meta = meta_r.json()

        if "google-apps.document" in meta.get("mimeType", ""):
            # Export as DOCX — preserves Heading 1/2/3 styles reliably across all doc types
            # HTML export can emit class-based styles instead of semantic <h2> tags for
            # documents originally uploaded as HTML, causing section splitting to fail.
            er = await client.get(
                f"https://www.googleapis.com/drive/v3/files/{doc_id}/export",
                headers={"Authorization": f"Bearer {token}"},
                params={"mimeType": "application/vnd.openxmlformats-officedocument.wordprocessingml.document"},
                timeout=20,
            )
            text = _docx_to_markdown(er.content) if er.is_success else ""
        else:
            dr = await client.get(
                f"https://www.googleapis.com/drive/v3/files/{doc_id}?alt=media",
                headers={"Authorization": f"Bearer {token}"},
                timeout=20,
            )
            text = extract_text(Path(meta["name"]).suffix.lower(), dr.content) if dr.is_success else ""

    if not text.strip():
        raise HTTPException(status_code=500, detail="Could not read document content")

    sections = _parse_sections(text)
    title_match = re.search(r'^#\s+(.+)$', text, re.MULTILINE)
    return {
        "title": title_match.group(1).strip() if title_match else meta["name"],
        "doc_type": doc_type,
        "drive_link": meta.get("webViewLink", ""),
        "sections": sections,
        "full_content": text,
    }


class DocRegenerateRequest(BaseModel):
    section_heading: str
    instruction: str
    doc_type: str
    full_content: str


def _get_doc_model(doc_type: str) -> genai.GenerativeModel:
    if doc_type == "brd":
        return get_brd_model(load_brd_context())
    if doc_type == "ird":
        return get_ird_model(load_ird_context())
    return get_model(load_context())


def _doc_field_key(doc_type: str) -> str:
    return "brd" if doc_type == "brd" else "crd"


async def _update_drive_doc(doc_id: str, markdown_content: str) -> None:
    import markdown as md_lib
    html_body = md_lib.markdown(markdown_content, extensions=["tables", "fenced_code"])
    html = f"<!DOCTYPE html><html><body>{html_body}</body></html>"

    token = await _get_drive_token(readonly=False)
    boundary = "boundary_sep_417a"
    body = (
        f"--{boundary}\r\nContent-Type: application/json\r\n\r\n{{}}\r\n"
        f"--{boundary}\r\nContent-Type: text/html\r\n\r\n{html}\r\n"
        f"--{boundary}--"
    ).encode("utf-8")

    async with httpx.AsyncClient(timeout=30) as client:
        r = await client.patch(
            f"https://www.googleapis.com/upload/drive/v3/files/{doc_id}",
            headers={
                "Authorization": f"Bearer {token}",
                "Content-Type": f"multipart/related; boundary={boundary}",
            },
            params={"uploadType": "multipart"},
            content=body,
        )
    if not r.is_success:
        logger.warning("Drive update failed for %s: %s", doc_id, r.text)
        raise HTTPException(status_code=502, detail=f"Failed to save document to Drive: {r.text}")


# ----- Surgical section replacement via the Google Docs API -----
# Instead of re-uploading the whole document (which makes Google Docs re-render
# everything and lose the original formatting), we replace ONLY the regenerated
# section's paragraphs in place, leaving every other section byte-for-byte intact.

def _u16len(s: str) -> int:
    """Length of a string in UTF-16 code units — what the Docs API uses for indices."""
    return len(s.encode("utf-16-le")) // 2


def _extract_bold(text: str) -> tuple[str, list[tuple[int, int]]]:
    """Strip **bold** markers, returning clean text + bold spans as UTF-16 offsets."""
    bolds: list[tuple[int, int]] = []
    out: list[str] = []
    pos = 0
    bold_on = False
    for part in re.split(r'(\*\*)', text):
        if part == '**':
            bold_on = not bold_on
            continue
        if not part:
            continue
        start = pos
        out.append(part)
        pos += _u16len(part)
        if bold_on:
            bolds.append((start, pos))
    return "".join(out), bolds


def _parse_markdown_blocks(section_md: str) -> list[dict]:
    """Parse a section's markdown into per-paragraph blocks for the Docs API.

    Each block: {kind: heading1..6|bullet|number|normal, text, bolds}.
    Mirrors how _docx_to_markdown emits one paragraph per line.
    """
    blocks: list[dict] = []
    for raw in section_md.split("\n"):
        line = raw.strip()
        if not line:
            continue
        kind = "normal"
        text = line
        m = re.match(r'^(#{1,6})\s+(.*)$', text)
        if m:
            kind = f"heading{len(m.group(1))}"
            text = m.group(2).strip()
        elif re.match(r'^[-*]\s+', text):
            kind = "bullet"
            text = re.sub(r'^[-*]\s+', '', text)
        elif re.match(r'^\d+\.\s+', text):
            kind = "number"
            text = re.sub(r'^\d+\.\s+', '', text)
        clean, bolds = _extract_bold(text)
        blocks.append({"kind": kind, "text": clean, "bolds": bolds})
    return blocks


async def _replace_doc_section(doc_id: str, section_heading: str, next_heading: str | None, new_section_md: str) -> bool:
    """Replace one section in a Google Doc in place via the Docs API.

    Returns True on success, False if the section's heading couldn't be located
    in the doc (e.g. the doc has no real heading styles) so the caller can fall back.
    """
    token = await _get_drive_token(readonly=False)
    headers = {"Authorization": f"Bearer {token}"}
    async with httpx.AsyncClient(timeout=30) as client:
        r = await client.get(f"https://docs.googleapis.com/v1/documents/{doc_id}", headers=headers)
        if not r.is_success:
            logger.warning("Docs API read failed for %s: %s", doc_id, r.text)
            return False
        content = r.json().get("body", {}).get("content", [])

        paras = []
        for el in content:
            p = el.get("paragraph")
            if not p:
                continue
            txt = "".join(e.get("textRun", {}).get("content", "") for e in p.get("elements", []))
            paras.append({"text": txt.strip("\n").strip(), "start": el["startIndex"], "end": el["endIndex"]})
        if not paras:
            return False
        body_end = content[-1]["endIndex"]

        target = section_heading.strip()
        si = next((i for i, p in enumerate(paras) if p["text"] == target), None)
        if si is None:
            return False
        section_start = paras[si]["start"]

        section_end = None
        if next_heading:
            nh = next_heading.strip()
            for p in paras[si + 1:]:
                if p["text"] == nh:
                    section_end = p["start"]
                    break
            if section_end is None:
                # The next heading was expected but couldn't be located in the live
                # doc. Treating this as the "last section" would delete everything
                # from here to the end of the document, so bail and let the caller
                # fall back to a safe full re-upload instead of losing content.
                logger.warning(
                    "next_heading %r not found in doc %s — aborting surgical replace to avoid data loss",
                    next_heading, doc_id,
                )
                return False
            is_last = False
        else:
            is_last = True
            section_end = body_end - 1  # genuinely the last section; keep final newline

        blocks = _parse_markdown_blocks(new_section_md)
        if not blocks:
            return False

        # Build the plain text to insert (one paragraph per block).
        if is_last:
            insert_text = "\n".join(b["text"] for b in blocks)
        else:
            insert_text = "".join(b["text"] + "\n" for b in blocks)

        requests: list[dict] = [
            {"deleteContentRange": {"range": {"startIndex": section_start, "endIndex": section_end}}},
            {"insertText": {"location": {"index": section_start}, "text": insert_text}},
        ]

        # Style each inserted paragraph. Inserted text inherits the heading style at
        # the insertion point, so we explicitly set the style of every paragraph.
        pos = section_start
        for b in blocks:
            tlen = _u16len(b["text"])
            p_start, p_end = pos, pos + tlen
            named = f"HEADING_{b['kind'][-1]}" if b["kind"].startswith("heading") else "NORMAL_TEXT"
            # Reset inherited bold across the paragraph, then re-apply real bold spans.
            requests.append({"updateTextStyle": {
                "range": {"startIndex": p_start, "endIndex": max(p_end, p_start + 1)},
                "textStyle": {"bold": False}, "fields": "bold"}})
            for bs, be in b["bolds"]:
                if be > bs:
                    requests.append({"updateTextStyle": {
                        "range": {"startIndex": p_start + bs, "endIndex": p_start + be},
                        "textStyle": {"bold": True}, "fields": "bold"}})
            requests.append({"updateParagraphStyle": {
                "range": {"startIndex": p_start, "endIndex": p_end + 1},
                "paragraphStyle": {"namedStyleType": named}, "fields": "namedStyleType"}})
            if b["kind"] in ("bullet", "number"):
                preset = "BULLET_DISC_CIRCLE_SQUARE" if b["kind"] == "bullet" else "NUMBERED_DECIMAL_ALPHA_ROMAN"
                requests.append({"createParagraphBullets": {
                    "range": {"startIndex": p_start, "endIndex": p_end + 1},
                    "bulletPreset": preset}})
            pos = p_end + 1

        br = await client.post(
            f"https://docs.googleapis.com/v1/documents/{doc_id}:batchUpdate",
            headers=headers,
            json={"requests": requests},
        )
        if not br.is_success:
            logger.warning("Docs API batchUpdate failed for %s: %s", doc_id, br.text)
            return False
    return True


@app.post("/documents/{doc_id}/regenerate-section")
async def regenerate_doc_section(doc_id: str, req: DocRegenerateRequest, _: dict = Depends(require_auth)):
    if not req.section_heading.strip():
        raise HTTPException(status_code=400, detail="section_heading is required")

    model = _get_doc_model(req.doc_type)
    field_key = _doc_field_key(req.doc_type)
    instruction_line = f"\nAdditional instructions: {req.instruction.strip()}" if req.instruction.strip() else ""

    response = await model.generate_content_async(
        f"""You are given a complete document in Markdown for context. Rewrite ONLY the one section named below and return JUST that section — not the whole document.

Section to regenerate: {req.section_heading.strip()}{instruction_line}

FORMATTING RULES — follow the document template in your corporate context exactly for this section:
- Start your output with the section's heading line, keeping the heading text identical to "{req.section_heading.strip()}" so it can be matched in the document.
- Use the same Markdown heading levels the template defines for this section and its sub-sections (e.g. `##` for a main section, `###` for sub-sections). Do NOT use bold text (`**...**`) in place of a heading.
- Use Markdown bullet lists (`- `) for any list of items, exactly as the template does.
- Match the template's structure, field labels, and ordering for this section.

Full document (context only — do NOT return the whole thing):
{req.full_content}

Return ONLY the rewritten section in Markdown, starting with its heading "{req.section_heading.strip()}". No other sections, no preamble, no explanation."""
    )
    new_section_md = strip_code_fences(response.text).strip()

    # Rebuild the full document by splicing the rewritten section into the
    # ORIGINAL content, preserving every other section (and the title preamble)
    # verbatim. We never trust the model to echo back the whole document — models
    # commonly return only the section, and reconstructing from the original is
    # what guarantees the rest of the document can't be lost when it's saved.
    target = req.section_heading.strip().lower()
    parts = re.split(r'(?=\n#{1,3} )', req.full_content)
    rebuilt: list[str] = []
    matched = False
    for part in parts:
        first_line = part.strip().split('\n')[0] if part.strip() else ''
        m = re.match(r'^#{1,3}\s+(.+)', first_line)
        heading = m.group(1).strip().lower() if m else None
        if not matched and heading == target:
            leading = part[:len(part) - len(part.lstrip('\n'))]  # keep original spacing
            rebuilt.append(leading + new_section_md)
            matched = True
        else:
            rebuilt.append(part)
    reconstructed_full = ''.join(rebuilt) if matched else req.full_content

    updated_section = {"heading": req.section_heading.strip(), "content": new_section_md}

    # Find the heading that follows this section (boundary for the surgical replace).
    orig_headings = [s["heading"] for s in _parse_sections(req.full_content)]
    next_heading = None
    if req.section_heading in orig_headings:
        idx = orig_headings.index(req.section_heading)
        if idx + 1 < len(orig_headings):
            next_heading = orig_headings[idx + 1]

    # This endpoint returns a PREVIEW only — it does NOT write to Drive. The
    # caller shows the proposed change for review, then calls
    # POST /documents/{doc_id}/save-section to persist it once the user confirms.
    return {
        field_key: reconstructed_full,
        "updated_section": updated_section,
        "full_content": reconstructed_full,
        "next_heading": next_heading,
    }


class DocSaveSectionRequest(BaseModel):
    doc_type: str = "crd"
    section_heading: str
    next_heading: str | None = None
    section_content: str
    full_content: str


async def _full_content_is_safe(doc_id: str, full_content: str) -> bool:
    """Guard for the full re-upload path: a section edit should never shrink the
    document to a fragment. Compares the incoming content's length against the
    live doc's text (format-agnostic, so it also protects old bold-heading docs).
    Fails open — if the live doc can't be read, the save is allowed."""
    new_len = len((full_content or "").strip())
    try:
        token = await _get_drive_token(readonly=True)
        async with httpx.AsyncClient(timeout=30) as client:
            r = await client.get(
                f"https://docs.googleapis.com/v1/documents/{doc_id}",
                headers={"Authorization": f"Bearer {token}"},
            )
        if not r.is_success:
            return True
        live_len = 0
        for el in r.json().get("body", {}).get("content", []):
            p = el.get("paragraph")
            if not p:
                continue
            live_len += len("".join(
                e.get("textRun", {}).get("content", "") for e in p.get("elements", [])
            ))
    except Exception as exc:
        logger.warning("save-section safety check failed for %s: %s", doc_id, exc)
        return True
    # Only block when the live doc has real substance and the incoming content is
    # less than half its size — the signature of a truncated single-section payload.
    if live_len > 800 and new_len < live_len * 0.5:
        return False
    return True


@app.post("/documents/{doc_id}/save-section")
async def save_doc_section(doc_id: str, req: DocSaveSectionRequest, _: dict = Depends(require_auth)):
    """Persist a previewed section edit to the Drive doc. Split out from
    regenerate-section so the user reviews and confirms the change before it's
    written. Surgical in-place edit, with a full re-upload fallback."""
    if not req.section_heading.strip():
        raise HTTPException(status_code=400, detail="section_heading is required")

    format_preserved = False
    if req.section_content.strip():
        try:
            format_preserved = await _replace_doc_section(
                doc_id, req.section_heading, req.next_heading, req.section_content
            )
        except Exception as exc:
            logger.warning("Surgical section replace failed for %s: %s", doc_id, exc)
            format_preserved = False

    if not format_preserved:
        # Full re-upload fallback. Guard against wiping the document: if the
        # incoming full_content is far shorter than the live doc (e.g. a stale
        # client sent only the regenerated section), refuse rather than overwrite
        # the whole doc with a fragment.
        if not await _full_content_is_safe(doc_id, req.full_content):
            logger.warning(
                "save-section: full_content for %s looks truncated — refusing full re-upload",
                doc_id,
            )
            raise HTTPException(
                status_code=409,
                detail="The update looked incomplete and was not saved, to avoid "
                       "overwriting the rest of the document. Please retry.",
            )
        await _update_drive_doc(doc_id, req.full_content)

    return {
        "full_content": req.full_content,
        "format_preserved": format_preserved,
    }


class ApplyHeadingStylesRequest(BaseModel):
    markdown: str


@app.post("/documents/{doc_id}/apply-heading-styles")
async def apply_doc_heading_styles(doc_id: str, req: ApplyHeadingStylesRequest, _: dict = Depends(require_auth)):
    """After a Drive HTML upload, apply real Heading 1/2/3 styles to paragraphs
    whose text matches headings in the source markdown.  Best-effort — errors
    are logged but not surfaced to the user since the doc is already uploaded."""
    heading_map: dict[str, int] = {}
    for line in req.markdown.splitlines():
        m = re.match(r'^(#{1,6})\s+(.+)$', line)
        if m:
            heading_map[m.group(2).strip()] = len(m.group(1))

    if not heading_map:
        return {"status": "ok", "updated": 0}

    token = await _get_drive_token(readonly=False)
    headers = {"Authorization": f"Bearer {token}"}
    async with httpx.AsyncClient(timeout=30) as client:
        r = await client.get(f"https://docs.googleapis.com/v1/documents/{doc_id}", headers=headers)
        if not r.is_success:
            logger.warning("apply-heading-styles: Docs read failed for %s: %s", doc_id, r.text)
            return {"status": "warning", "message": "could not read document"}

        body_content = r.json().get("body", {}).get("content", [])
        requests_list: list[dict] = []
        for el in body_content:
            p = el.get("paragraph")
            if not p:
                continue
            txt = "".join(
                e.get("textRun", {}).get("content", "") for e in p.get("elements", [])
            ).strip("\n").strip()
            if txt in heading_map:
                level = heading_map[txt]
                requests_list.append({"updateParagraphStyle": {
                    "range": {"startIndex": el["startIndex"], "endIndex": el["endIndex"]},
                    "paragraphStyle": {"namedStyleType": f"HEADING_{level}"},
                    "fields": "namedStyleType",
                }})

        if not requests_list:
            return {"status": "ok", "updated": 0}

        br = await client.post(
            f"https://docs.googleapis.com/v1/documents/{doc_id}:batchUpdate",
            headers=headers,
            json={"requests": requests_list},
        )
        if not br.is_success:
            logger.warning("apply-heading-styles: batchUpdate failed for %s: %s", doc_id, br.text)
            return {"status": "warning", "message": br.text}

    return {"status": "ok", "updated": len(requests_list)}


@app.post("/brd/log-to-sheet")
async def brd_log_to_sheet(req: BrdLogToSheetRequest, _: dict = Depends(require_auth)):
    date_str = datetime.date.today().strftime("%-d %b %y")
    title = extract_brd_title(req.brd)
    objective = extract_brd_objective(req.brd)
    crds = extract_brd_crds(req.brd)
    initials = generate_br_initials(title)

    def _write():
        ws = _open_worksheet("BR")

        existing_ids = ws.col_values(1)
        prefix = f"{initials}-"
        nums = []
        for id_val in existing_ids:
            s = str(id_val)
            if s.startswith(prefix):
                try:
                    nums.append(int(s[len(prefix):]))
                except ValueError:
                    pass
        br_id = f"{initials}-{max(nums, default=0) + 1}"

        # Column order (BR worksheet, 2026 catalogue):
        #   1 ID | 2 Business Requirement | 3 Business Objective | 4 CRDID
        #   5 IRDID | 6 Priority | 7 Date Submitted | 8 BRD Document
        #   9 FeatureIDs | 10 Feature Names | 11 Release | 12 End to End
        # IRDID/Priority and cols 9–12 aren't in the BRD doc, so they're left
        # blank for manual entry.
        ws.append_row(
            [br_id, title, objective, crds, "", "", date_str, req.drive_link],
            value_input_option="RAW",
            insert_data_option="INSERT_ROWS",
        )

    try:
        await asyncio.to_thread(_write)
    except Exception as exc:
        logger.warning("BR sheet write failed in /brd/log-to-sheet (%s)", exc)
        return {"status": "warning", "message": str(exc)}
    return {"status": "ok"}
